import time
import os
import sys
from datetime import datetime
from zoneinfo import ZoneInfo
from io import BytesIO
from urllib.parse import urlparse, parse_qs

from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.common.by import By

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


# ================= CONFIG =================

captured_images = []
TEMPLATE_FILE = "template_grafana.pptx"

GRAFANA_USER = ""
GRAFANA_PASS = ""


# ================= ARGUMENT =================

if len(sys.argv) != 5:
    print("Usage:")
    print("python weekly_capture_grafana.HO.py START_DATE START_TIME END_DATE END_TIME")
    sys.exit(1)

START_DATE = sys.argv[1]
START_TIME = sys.argv[2]
END_DATE = sys.argv[3]
END_TIME = sys.argv[4]


# ================= TIME =================

def to_grafana_ts(date_str, time_str):
    dt = datetime.strptime(f"{date_str} {time_str}", "%Y-%m-%d %H:%M")
    dt = dt.replace(tzinfo=ZoneInfo("Asia/Jakarta"))
    return int(dt.timestamp() * 1000)

FROM_TS = to_grafana_ts(START_DATE, START_TIME)
TO_TS = to_grafana_ts(END_DATE, END_TIME)


# ================= HOSTNAME =================

def get_hostname_from_url(url):
    parsed = urlparse(url)
    params = parse_qs(parsed.query)

    keys = ["var-nodename","var-instance_name","var-instance","var-vm"]

    for key in keys:
        if key in params:
            val = params[key][0]
            if ":" in val:
                val = val.split(":")[0]
            if val.lower() != "all":
                return val

    return "dashboard"


# ================= CHROME =================

chrome_options = Options()
chrome_options.add_argument("--headless=new")
chrome_options.add_argument("--no-sandbox")
chrome_options.add_argument("--disable-dev-shm-usage")

driver = webdriver.Chrome(options=chrome_options)
wait = WebDriverWait(driver, 60)

driver.set_window_size(1920, 1080)


# ================= LOGIN =================

print("Login Grafana...")

driver.get("https://monitoring.ilcs.co.id/login")

wait.until(lambda d: d.find_element(By.NAME, "user"))

driver.find_element(By.NAME, "user").send_keys(GRAFANA_USER)
driver.find_element(By.NAME, "password").send_keys(GRAFANA_PASS)
driver.find_element(By.CSS_SELECTOR, "button[type='submit']").click()

time.sleep(8)

print("Login berhasil")


# ================= URL =================

URL_LIST = [

f"https://monitoring.ilcs.co.id/d/df07rysvmjx1cc/vm-linux?orgId=1&var-prome_ds=eey4rwzgb9ibke&var-job=node-exporter&var-nodename=gcp-srvlpmy01&var-instance=10.174.3.60:9100&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/df07rysvmjx1cc/vm-linux?orgId=1&var-prome_ds=eey4rwzgb9ibke&var-job=node-exporter&var-nodename=gcp-srvlppajak01&var-instance=10.174.3.62:9100&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/df07rysvmjx1cc/vm-linux?orgId=1&var-prome_ds=eey4rwzgb9ibke&var-job=node-exporter&var-nodename=gcp-srvlppeproc01&var-instance=10.174.3.3:9100&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/df07rysvmjx1cc/vm-linux?orgId=1&var-prome_ds=eey4rwzgb9ibke&var-job=node-exporter&var-nodename=gcp-srvlpprima01&var-instance=10.174.3.53:9100&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/df07rysvmjx1cc/vm-linux?orgId=1&var-prome_ds=eey4rwzgb9ibke&var-job=node-exporter&var-nodename=gcp-srvlprkm01&var-instance=10.174.3.4:9100&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/bey7rgxvk9b0ga/vm-windows?orgId=1&var-prome_ds=eey4rwzgb9ibke&var-nodename=gcp-srvwpbios01&var-instance=10.174.3.6:9182&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/bey7rgxvk9b0ga/vm-windows?orgId=1&var-prome_ds=eey4rwzgb9ibke&var-nodename=gcp-srvwpgatewaypowerbi01&var-instance=10.174.3.5:9182&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/bey7rgxvk9b0ga/vm-windows?orgId=1&var-prome_ds=eey4rwzgb9ibke&var-nodename=gcp-srvwpsimtax&var-instance=10.174.3.61:9182&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/df07rysvmjx1cc/vm-linux?orgId=1&var-prome_ds=eey4rwzgb9ibke&var-job=node-exporter&var-nodename=gcp-srvlpcentra01&var-instance=10.174.3.199:9100&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/df07rysvmjx1cc/vm-linux?orgId=1&var-prome_ds=eey4rwzgb9ibke&var-job=node-exporter&var-nodename=gcp-svrlptanos01&var-instance=10.174.3.204:9100&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/ferdmagigyzggb/k8s-back-office?orgId=1&var-datasource=cecpye14uwxkwc&var-project=pelindo-back-office&var-cluster_name=back-office-new&var-namespace=&var-container_name=&var-alignmentPeriod=grafana-auto&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/ferdmagigyzggb/k8s-back-office?orgId=1&var-datasource=cecpye14uwxkwc&var-project=pelindo-back-office&var-cluster_name=back-office-2&var-namespace=&var-container_name=&var-alignmentPeriod=grafana-auto&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/df07rysvmjx1cc/vm-linux?orgId=1&var-prome_ds=eey4rwzgb9ibke&var-job=node-exporter&var-nodename=ematerai&var-instance=10.174.2.6:9100&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/bey7rgxvk9b0ga/vm-windows?orgId=1&var-prome_ds=eey4rwzgb9ibke&var-nodename=ibs-api-1&var-instance=10.174.2.40:9182&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/bey7rgxvk9b0ga/vm-windows?orgId=1&var-prome_ds=eey4rwzgb9ibke&var-nodename=ibs-api-2&var-instance=10.174.2.42:9182&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/bey7rgxvk9b0ga/vm-windows?orgId=1&var-prome_ds=eey4rwzgb9ibke&var-nodename=integration&var-instance=10.174.2.7:9182&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/advadwnf34b34/vm-payment?orgId=1&var-datasource=cecpye14uwxkwc&var-project=pelindo-payment&var-instance_name=pelindo-pay&var-alignmentPeriod=grafana-auto&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/bey7rgxvk9b0ga/vm-windows?orgId=1&var-prome_ds=eey4rwzgb9ibke&var-nodename=payment-bank-01&var-instance=10.174.2.50:9182&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/87v4vm98cfrvr2v3/k8s-payment?orgId=1&var-datasource=cecpye14uwxkwc&var-project=pelindo-payment&var-cluster_name=k8s-payment&var-container_name=&var-alignmentPeriod=grafana-auto&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/advadwnf34b34/vm-payment?orgId=1&var-datasource=cecpye14uwxkwc&var-project=pelindo-payment&var-instance_name=All&var-alignmentPeriod=grafana-auto&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/ceu58oumv3qiod/k8s-sap?orgId=1&var-datasource=cecpye14uwxkwc&var-project=pelindo-sap&var-cluster_name=sap-peic&var-namespace=&var-container_name=&var-alignmentPeriod=grafana-auto&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/aeu5916cfad4wb/sql-sap?orgId=1&var-datasource=cecpye14uwxkwc&var-project=pelindo-sap&var-db_name=svreicprdsql&var-alignmentPeriod=grafana-auto&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/eeu5c3dxxotfkc/redis-sap?orgId=1&var-datasource=cecpye14uwxkwc&var-project=pelindo-sap&var-instance_name=All&var-alignmentPeriod=grafana-auto&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/den1cqwgkwa9sa/vm-qa?orgId=1&var-datasource=cecpye14uwxkwc&var-project=pelindo-qa&var-instance_name=All&var-alignmentPeriod=grafana-auto&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/eexfbo85ofwu8e/interconnect?orgId=1&from=1772977326045&to=1772988126045&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/cey479vts11xce/gcp-load-balancer?orgId=1&from=1772966586076&to=1772988186076&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/cefogwvbsp2bkc/cloud-armor-all?orgId=1&var-cloud_armor_name=All&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/be954464vl2pseasc/vm-apps?orgId=1&var-instance_name=oci-srvlptravel01&var-interval=1m&var-compartment=pelindo%20%3E%20App-Pelindo-HO&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/be954464vl2pse/vm-cloudera?orgId=1&var-instance_name=dlake-repo&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/be954464vl2pse/vm-cloudera?orgId=1&var-instance_name=dlake-edge&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/be954464vl2pse/vm-cloudera?orgId=1&var-instance_name=dlake-wkr&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/be954464vl2pse/vm-cloudera?orgId=1&var-instance_name=dlake-mstr&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/be954464vl2pse/vm-cloudera?orgId=1&var-instance_name=dlake-util&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/be954464vl2pse/vm-cloudera?orgId=1&var-instance_name=strm-kfk&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/be954464vl2pse/vm-cloudera?orgId=1&var-instance_name=strm-nf&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/be954464vl2pse/vm-cloudera?orgId=1&var-instance_name=strm-mstr&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/be954464vl2pse/vm-cloudera?orgId=1&var-instance_name=strm-util&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/be954464vl2pse/vm-cloudera?orgId=1&var-instance_name=freeipa&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/be954464vl2pse/vm-cloudera?orgId=1&var-instance_name=dns-serve&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/fegsg889m1vk0d/vm-phinnisi?orgId=1&var-instance_name=All&var-compartment=ilcspsdcloudaccount%20%3E%20prod-phinnisi&var-interval=1m&var-nodepool1=oke-c5r74d6mfva-ndgzllb7uca-s2dlrd4cwu&var-nodepool2=oke-c5r74d6mfva-nlflfmfum5q-s2dlrd4cwua&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/deiv4bn1a5w5ce/db-phinnisi?orgId=1&var-db_name=All&var-compartment=ilcspsdcloudaccount%20%3E%20prod-phinnisi&var-interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/y3gun4vh5ycefng78/db-oracle-db-pelindo?orgId=1&var-db_name=madra&var-interval=1m&var-compartment=pelindo%20%3E%20DB-Pelindo-HO&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/y3gun4vh5ycefng78/db-oracle-db-pelindo?orgId=1&var-db_name=pancala&var-interval=1m&var-compartment=pelindo%20%3E%20DB-Pelindo-HO&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/y3gun4vh5ycefng78/db-oracle-db-pelindo?orgId=1&var-db_name=hastina&var-interval=1m&var-compartment=pelindo%20%3E%20DB-Pelindo-HO&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/y3gun4vh5ycefng78/db-oracle-db-pelindo?orgId=1&var-db_name=ayodya&var-interval=1m&var-compartment=pelindo%20%3E%20DB-Pelindo-HO&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/y3gun4vh5ycefng78/db-oracle-db-pelindo?orgId=1&var-db_name=alengka&var-interval=1m&var-compartment=pelindo%20%3E%20DB-Pelindo-HO&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/y3gun4vh5yceftm6u/db-mysql-db-pelindo?orgId=1&var-db_name=All&var-interval=1m&var-compartment=pelindo%20%3E%20DB-Pelindo-HO&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/cejrks918xtkwb/db-sigap-postgres?orgId=1&var-dbname=sigap-prod&var-Interval=1m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/cezp5v5ej4f0gc/centralized-kubenertes?orgId=1&var-interval=30m&var-cluster=huwai-portaverse&var-namespace=All&var-container=All&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/cdyyfytqtixogfdsfdsf/centralized-rds?orgId=1&var-EPS=0&var-hostIP=10.95.200.227&var-hostname=dbportaverse&var-total=4&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/portefaix_node_fleet_overviewdsd/centralized-vm?orgId=1&refresh=30s&var-DS_Metrics=eey51w20ssn40e&var-job=All&var-vm=simon-stid&var-rate_interval=5m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/cdyyfytqtixogfdsfdsf/centralized-rds?orgId=1&var-EPS=0&var-hostIP=10.95.200.49&var-hostname=db-stid-simon&var-total=4&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/portefaix_node_fleet_overviewdsd/centralized-vm?orgId=1&refresh=30s&var-DS_Metrics=eey51w20ssn40e&var-job=All&var-vm=ims&var-rate_interval=5m&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

f"https://monitoring.ilcs.co.id/d/cdyyfytqtixogfdsfdsf/centralized-rds?orgId=1&var-EPS=0&var-hostIP=10.95.200.11&var-hostname=db-apps&var-total=4&from={FROM_TS}&to={TO_TS}&timezone=Asia%2FJakarta",

]


# ================= CAPTURE =================

for url in URL_LIST:

    name = get_hostname_from_url(url)
    print("Opening:", name)

    driver.get(url)

    wait.until(lambda d: d.execute_script("return document.readyState") == "complete")

    time.sleep(10)

    driver.execute_script("window.scrollTo(0,0)")
    time.sleep(2)

    driver.execute_script("document.body.style.zoom='90%'")
    time.sleep(2)

    png = driver.get_screenshot_as_png()
    image = Image.open(BytesIO(png))

    width, height = image.size

    image = image.crop((0, 0, width, int(height * 0.87)))

    captured_images.append((name, image))

    print("Captured:", name)


# ================= PPT =================

print("Generating PPT...")

ppt_name = f"Grafana_Report_{START_DATE}_to_{END_DATE}.pptx"

if os.path.exists(ppt_name):
    os.remove(ppt_name)

prs = Presentation(TEMPLATE_FILE)

slide_mapping = {
    0: 3,    # gambar1 -> slide4
    1: 5,    # gambar2 -> slide6
    2: 7,    # gambar3 -> slide8
    3: 9,    # gambar4 -> slide10
    4: 11,   # gambar5 -> slide12
    5: 13,   # gambar6 -> slide14
    6: 15,   # gambar7 -> slide16
    7: 17,   # gambar8 -> slide18
    8: 19,   # gambar9 -> slide20
    9: 21,   # gambar10 -> slide22
    10: 23,  # gambar11 -> slide24
    11: 25,  # gambar12 -> slide26
    12: 28,  # gambar13 -> slide29
    13: 30,  # gambar14 -> slide31
    14: 31,  # gambar15 -> slide32
    15: 33,  # gambar16 -> slide34
    16: 35,  # gambar17 -> slide36
    17: 37,  # gambar18 -> slide38
    18: 39,  # gambar19 -> slide40
    19: 42,  # gambar20 -> slide43
    20: 44,  # gambar21 -> slide45
    21: 46,  # gambar22 -> slide47
    22: 48,  # gambar23 -> slide49
    23: 51,  # gambar24 -> slide52
    24: 53,  # gambar25 -> slide54
    25: 54,  # gambar26 -> slide55
    26: 55,  # gambar27 -> slide56
    27: 58,  # gambar28 -> slide59
    28: 62,  # gambar29 -> slide63
    29: 63,  # gambar30 -> slide64
    30: 64,  # gambar31 -> slide65
    31: 65,  # gambar32 -> slide66
    32: 66,  # gambar33 -> slide67
    33: 67,  # gambar34 -> slide68
    34: 68,  # gambar35 -> slide69
    35: 69,  # gambar36 -> slide70
    36: 70,  # gambar37 -> slide71
    37: 71,  # gambar38 -> slide72
    38: 72,  # gambar39 -> slide73
    39: 74,  # gambar40 -> slide75
    40: 76,  # gambar41 -> slide77
    41: 78,  # gambar42 -> slide79
    42: 80,  # gambar43 -> slide81
    43: 82,  # gambar44 -> slide83
    44: 84,  # gambar45 -> slide85
    45: 86,  # gambar46 -> slide87
    46: 88,  # gambar47 -> slide89
    47: 90,  # gambar48 -> slide91
    48: 93,  # gambar49 -> slide94
    49: 95,  # gambar50 -> slide96
    50: 98,  # gambar51 -> slide99
    51: 100, # gambar52 -> slide101
    52: 102, # gambar53 -> slide103
    53: 104  # gambar54 -> slide105
}

for idx, (name, image) in enumerate(captured_images):

    if idx not in slide_mapping:
        continue

    slide = prs.slides[slide_mapping[idx]]

    if slide.shapes.title:
        slide.shapes.title.text = f"{name} ({START_DATE} to {END_DATE})"

    img_stream = BytesIO()
    image.save(img_stream, format="PNG")
    img_stream.seek(0)

    slide.shapes.add_picture(
        img_stream,
        Inches(0.3),
        Inches(1.3),
        width=Inches(12.8)
    )

prs.save(ppt_name)

print("DONE:", ppt_name)
